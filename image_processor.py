import os
import logging
import json
import hashlib
from PIL import Image
from pdf2image import convert_from_path
from PyPDF2 import PdfMerger
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH
from typing import List, Dict, Any
import tempfile
from datetime import datetime
import cv2
import numpy as np
import time

# PPTX imports
try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches, Pt as PPTPt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# PDF generation imports
try:
    from fpdf import FPDF
    FPDF_AVAILABLE = True
except ImportError:
    FPDF_AVAILABLE = False

import base64

# Gemini AI imports
try:
    from google import genai
    from PIL import Image as PILImage
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    logging.warning("google-genai not available.")

# Anthropic AI imports
try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False
    logging.warning("anthropic library not available.")

# Load configuration
try:
    import config
    from config import GEMINI_API_KEY, GEMINI_MODEL, ENABLE_AI_INSIGHTS
    print(f"[IMAGE_PROCESSOR] Loaded config OK. GEMINI_MODEL={GEMINI_MODEL}")
    logging.info(f"AI CONFIG: Primary Model={GEMINI_MODEL}, Provider={getattr(config, 'AI_PROVIDER', 'gemini')}")
except Exception as _cfg_err:
    import os
    GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
    GEMINI_MODEL = "gemini-2.5-flash"
    ENABLE_AI_INSIGHTS = False
    print(f"[IMAGE_PROCESSOR] Config import FAILED: {_cfg_err}. Using fallback model.")
    logging.warning(f"config.py not found or incomplete. Error: {_cfg_err}")

class ImageProcessor:
    def __init__(self):
        self.temp_files = []
        
        # Initialize Gemini AI client
        self.gemini_client = None
        self.anthropic_client = None
        
        if config.ENABLE_AI_INSIGHTS:
            if config.AI_PROVIDER == "anthropic" and ANTHROPIC_AVAILABLE:
                try:
                    self.anthropic_client = anthropic.Anthropic(api_key=config.ANTHROPIC_API_KEY)
                    logging.info("Anthropic client initialized successfully")
                except Exception as e:
                    logging.error(f"Failed to initialize Anthropic client: {e}")
            
            # If anthropic was not chosen or failed/unavailable, try Gemini
            if not self.anthropic_client and config.GEMINI_API_KEY and GEMINI_AVAILABLE:
                try:
                    self.gemini_client = genai.Client(api_key=config.GEMINI_API_KEY)
                    logging.info("Gemini client initialized successfully")
                except Exception as e:
                    logging.error(f"Failed to initialize Gemini client: {e}")
        else:
            reasons = []
            if not ENABLE_AI_INSIGHTS: reasons.append("ENABLE_AI_INSIGHTS is False")
            if not GEMINI_API_KEY: reasons.append("GEMINI_API_KEY is missing")
            if not GEMINI_AVAILABLE: reasons.append("google-genai library missing")
            logging.warning(f"Gemini AI is disabled. Reasons: {', '.join(reasons)}")

    def _format_filters(self, filters_dict: Dict) -> str:
        """Helper to format applied filters into a readable string"""
        if not filters_dict:
            return ""
        filter_strings = []
        for name, value in filters_dict.items():
            filter_strings.append(f"{name}: {value}")
        return ", ".join(filter_strings)
    
    def _format_datasource_info(self, datasources: List) -> str:
        """Helper to format datasource info into a readable string"""
        if not datasources:
            return ""
        parts = []
        for ds in datasources:
            name = ds.get('name', 'Data Source')
            has_extract = ds.get('hasExtracts', False)
            updated_at = ds.get('updatedAt')
            
            if updated_at:
                try:
                    from datetime import datetime as dt
                    date = dt.fromisoformat(updated_at.replace('Z', '+00:00'))
                    # Convert UTC to local system time
                    local_date = date.astimezone()
                    formatted_date = local_date.strftime('%b %d, %Y %I:%M %p')
                except:
                    formatted_date = updated_at
                
                if has_extract:
                    parts.append(f"Extract \u2014 As of {formatted_date}")
                else:
                    parts.append(f"Updated: {formatted_date}")
            elif has_extract:
                parts.append("Extract (refresh date unavailable)")
            else:
                parts.append("Live connection")
        
        return "; ".join(parts)
    
    def pdf_to_png(self, pdf_path: str, dpi: int = 200) -> str:
        """Convert PDF to PNG image"""
        try:
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=dpi)
            
            if not images:
                raise Exception("No images found in PDF")
            
            # Use the first page
            image = images[0]
            
            # Generate PNG filename
            base_name = os.path.splitext(os.path.basename(pdf_path))[0]
            png_path = os.path.join(os.path.dirname(pdf_path), f"{base_name}.png")
            
            # Save as PNG
            image.save(png_path, "PNG")
            
            logging.info(f"Successfully converted PDF to PNG: {png_path}")
            return png_path
            
        except Exception as e:
            logging.error(f"Failed to convert PDF to PNG: {str(e)}")
            raise Exception(f"PDF conversion failed: {str(e)}")
    
    def trim_to_dashboard_size(self, png_path: str, design_w: int = 0, design_h: int = 0) -> str:
        """
        Trim the blank whitespace margins that Tableau adds above and below the
        dashboard content when exporting to PDF (US Letter page).

        Uses pixel-level detection to find the actual first and last rows that
        contain non-background content, so no content is ever clipped regardless
        of how Tableau positions the content on the page.

        A small safety buffer is kept around the detected content boundary so that
        faint borders or anti-aliased edges are never lost.

        Returns the path to the trimmed PNG (saved alongside the original), or the
        original path unchanged if no meaningful whitespace is found.
        """
        BACKGROUND_THRESHOLD = 245   # pixels brighter than this are considered background
        SAFETY_BUFFER = 5            # extra pixels to keep above/below detected content
        MIN_WHITESPACE = 20          # only trim when at least this many blank rows exist

        try:
            img = Image.open(png_path)
            page_w, page_h = img.size

            # Work in grayscale for fast per-row min detection
            gray = img.convert('L')

            # Scan top-to-bottom for first row with any non-background pixel
            first_content_row = None
            for y in range(page_h):
                row = gray.crop((0, y, page_w, y + 1))
                if min(row.getdata()) < BACKGROUND_THRESHOLD:
                    first_content_row = y
                    break

            if first_content_row is None:
                logging.info("trim_to_dashboard_size: no content found, skipping trim.")
                return png_path

            # Scan bottom-to-top for last row with any non-background pixel
            last_content_row = first_content_row
            for y in range(page_h - 1, first_content_row, -1):
                row = gray.crop((0, y, page_w, y + 1))
                if min(row.getdata()) < BACKGROUND_THRESHOLD:
                    last_content_row = y
                    break

            top_whitespace    = first_content_row
            bottom_whitespace = page_h - last_content_row - 1

            logging.info(
                f"trim_to_dashboard_size: detected content rows {first_content_row}–{last_content_row} "
                f"(top_ws={top_whitespace}px, bot_ws={bottom_whitespace}px)"
            )

            # Only trim if there is meaningful whitespace on at least one side
            if top_whitespace < MIN_WHITESPACE and bottom_whitespace < MIN_WHITESPACE:
                logging.info("trim_to_dashboard_size: whitespace too small, skipping trim.")
                return png_path

            y_start = max(0,       first_content_row - SAFETY_BUFFER)
            y_end   = min(page_h,  last_content_row  + SAFETY_BUFFER + 1)

            logging.info(f"Trimming PNG to y=[{y_start}, {y_end}] (content + {SAFETY_BUFFER}px buffer)")

            cropped = img.crop((0, y_start, page_w, y_end))

            base, ext = os.path.splitext(png_path)
            trimmed_path = f"{base}_trimmed{ext}"
            cropped.save(trimmed_path, "PNG")
            logging.info(f"Trimmed PNG saved: {trimmed_path} ({cropped.size[0]}×{cropped.size[1]})")
            return trimmed_path

        except Exception as e:
            logging.error(f"trim_to_dashboard_size failed: {e}")
            return png_path  # fall back to original on any error

    def crop_image(self, image_path: str, crop_data: Dict[str, float]) -> str:
        """Crop an image based on crop coordinates"""
        try:
            image = Image.open(image_path)
            
            # Extract crop coordinates
            x1 = int(crop_data['x'])
            y1 = int(crop_data['y'])
            x2 = int(crop_data['x'] + crop_data['width'])
            y2 = int(crop_data['y'] + crop_data['height'])
            
            # Ensure coordinates are within image bounds
            x1 = max(0, min(x1, image.width))
            y1 = max(0, min(y1, image.height))
            x2 = max(0, min(x2, image.width))
            y2 = max(0, min(y2, image.height))
            
            # Ensure we have a valid crop area
            if x2 <= x1 or y2 <= y1:
                raise Exception("Invalid crop coordinates")
            
            # Crop the image
            cropped_image = image.crop((x1, y1, x2, y2))
            
            # Generate unique cropped filename with timestamp so successive crops
            # produce different filenames and the browser doesn't serve a cached version.
            import time as _time
            base_name = os.path.splitext(os.path.basename(image_path))[0]
            # Strip any previous _cropped_<ts> suffix so base stays clean
            if '_cropped' in base_name:
                base_name = base_name[:base_name.index('_cropped')]
            ts = int(_time.time() * 1000)
            cropped_path = os.path.join(os.path.dirname(image_path), f"{base_name}_cropped_{ts}.png")
            
            # Save cropped image
            cropped_image.save(cropped_path, "PNG")
            
            logging.info(f"Successfully cropped image: {cropped_path}")
            return cropped_path
            
        except Exception as e:
            logging.error(f"Failed to crop image: {str(e)}")
            raise Exception(f"Image cropping failed: {str(e)}")
    
    def combine_to_pdf(self, image_paths: List[str], output_dir: str, filename: str) -> str:
        """Combine multiple images into a single PDF (Legacy version, no text)"""
        try:
            output_path = os.path.join(output_dir, f"{filename}.pdf")
            
            # Convert images to PDF
            temp_pdfs = []
            merger = PdfMerger()
            
            for i, image_path in enumerate(image_paths):
                if not os.path.exists(image_path):
                    logging.warning(f"Image not found: {image_path}")
                    continue
                
                # Open and convert image to RGB if necessary
                image = Image.open(image_path)
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Create temporary PDF for this image
                temp_pdf_path = os.path.join(output_dir, f"temp_{i}.pdf")
                image.save(temp_pdf_path, "PDF")
                temp_pdfs.append(temp_pdf_path)
                
                # Add to merger
                merger.append(temp_pdf_path)
            
            if not temp_pdfs:
                raise Exception("No valid images to combine")
            
            # Write combined PDF
            merger.write(output_path)
            merger.close()
            
            # Clean up temporary PDFs
            for temp_pdf in temp_pdfs:
                try:
                    os.remove(temp_pdf)
                except:
                    pass
            
            logging.info(f"Successfully created combined PDF (raw): {output_path}")
            return output_path
            
        except Exception as e:
            logging.error(f"Failed to combine images to PDF: {str(e)}")
            raise Exception(f"PDF combination failed: {str(e)}")

    @staticmethod
    def _pdf_safe(text: str) -> str:
        """fpdf core fonts are latin-1 only — replace common unicode, drop the rest."""
        if not text:
            return ""
        repl = {
            "—": "-", "–": "-", "‘": "'", "’": "'",
            "“": '"', "”": '"', "•": "-", "·": "-",
            "▲": "+", "▼": "-", "↑": "+", "↓": "-",
            "→": "->", "≤": "<=", "≥": ">=", "✓": "OK",
            "…": "...",
        }
        for k, v in repl.items():
            text = text.replace(k, v)
        return text.encode("latin-1", "replace").decode("latin-1")

    def _build_report_dashboards(self, image_paths: List[str], summary_data: List[Dict]) -> List[Dict]:
        """Generate + parse AI insights for every dashboard up front and order
        the result worst-first (RED -> AMBER -> GREEN -> unrated). Shared by
        the executive PPTX and PDF builders."""
        dashboards = []
        for idx, image_path in enumerate(image_paths):
            if not os.path.exists(image_path):
                continue
            sd = summary_data[idx] if summary_data and idx < len(summary_data) else {}
            dashboard_name = sd.get("dashboard", f"Dashboard {idx + 1}")
            try:
                insights = self._generate_ai_insights(
                    image_path, dashboard_name, csv_data=sd.get("csv_data", ""))
                logging.info(f"Generated {len(insights)} insights for '{dashboard_name}'")
            except Exception as ai_e:
                logging.error(f"AI insights failed for '{dashboard_name}': {ai_e}")
                insights = []
            headline, sections = self._parse_insight_sections(insights)
            recommendation = next(
                (t for l, t in sections if l.upper().startswith("RECOMMENDATION")), "")
            status = next(
                (t.strip().upper() for l, t in sections if l.upper().startswith("STATUS")), "")
            if status not in ("RED", "AMBER", "GREEN"):
                status = ""
            key_metric = next(
                (t for l, t in sections if l.upper().startswith("KEY METRIC")), "")
            findings = [(l, t) for l, t in sections
                        if not l.upper().startswith(("RECOMMENDATION", "STATUS", "KEY METRIC"))]
            dashboards.append({
                "image_path": image_path,
                "name": dashboard_name,
                "filters": sd.get("applied_filters", {}),
                "datasources": sd.get("datasources", []),
                "headline": headline or dashboard_name,
                "findings": findings,
                "recommendation": recommendation,
                "status": status,
                "key_metric": key_metric,
            })
        # Management by exception: worst first
        status_rank = {"RED": 0, "AMBER": 1, "GREEN": 2}
        dashboards.sort(key=lambda d: status_rank.get(d["status"], 3))
        return dashboards

    def combine_to_pdf_with_details(self, image_paths: List[str], output_dir: str, filename: str, summary_data: List[Dict]) -> str:
        """Executive weekly-business-review PDF — same structure as the PPTX:
        brand band, numbered black section bars, executive scorecard with RAG
        statuses, yellow HEADLINE strips, image + comments, ACTION rows."""
        if not FPDF_AVAILABLE:
            logging.warning("fpdf2 not available, falling back to basic PDF")
            return self.combine_to_pdf(image_paths, output_dir, filename)

        import math

        try:
            output_path = os.path.join(output_dir, f"{filename}.pdf")
            dashboards = self._build_report_dashboards(image_paths, summary_data)

            W, H = 11.0, 8.5
            MX = 0.30
            CW = W - 2 * MX
            TOP, BOTTOM = 0.74, 8.02
            INK, MUTED = (26, 26, 26), (89, 89, 89)
            BAR, BAND, HDRBG = (17, 17, 17), (239, 239, 239), (217, 217, 217)
            LINEC, ROWALT, YELLOW = (175, 175, 175), (245, 245, 245), (255, 242, 0)
            RED, AMBER, GREEN = (192, 0, 0), (232, 156, 0), (46, 125, 50)
            REDBG, DARKRED = (251, 227, 227), (139, 0, 0)

            pdf = FPDF(orientation='L', unit='in', format='Letter')
            pdf.set_auto_page_break(auto=False)
            pdf.set_margins(MX, TOP, MX)
            page_no = 0

            def status_style(status):
                if status == "RED":
                    return RED, (255, 255, 255), "RED"
                if status == "AMBER":
                    return AMBER, INK, "AMBER"
                if status == "GREEN":
                    return GREEN, (255, 255, 255), "GREEN"
                return HDRBG, MUTED, "-"

            def new_page():
                nonlocal page_no
                page_no += 1
                pdf.add_page()
                pdf.set_fill_color(*BAND)
                pdf.rect(0, 0, W, 0.60, 'F')
                pdf.set_draw_color(*LINEC)
                pdf.set_line_width(0.01)
                pdf.line(0, 0.60, W, 0.60)
                pdf.set_text_color(*INK)
                pdf.set_font("helvetica", "B", 14)
                pdf.set_xy(0, 0.10)
                pdf.cell(W, 0.22, "BLEND360", align="C")
                pdf.set_text_color(*MUTED)
                pdf.set_font("helvetica", "", 7.5)
                pdf.set_xy(0, 0.345)
                pdf.cell(W, 0.14, "Dashboard Insights - Automated Report to Business Stakeholders", align="C")
                pdf.line(MX, 8.14, W - MX, 8.14)
                pdf.set_xy(MX, 8.18)
                pdf.set_font("helvetica", "", 6.5)
                pdf.cell(5.0, 0.12, self._pdf_safe(
                    datetime.now().strftime("Confidential  |  Generated %B %d, %Y")), align="L")
                pdf.set_xy(W - MX - 1.0, 8.18)
                pdf.cell(1.0, 0.12, f"Page {page_no}", align="R")
                return TOP

            def section_bar(y, label, status=None):
                pdf.set_fill_color(*BAR)
                pdf.rect(MX, y, CW, 0.26, 'F')
                pdf.set_text_color(255, 255, 255)
                pdf.set_font("helvetica", "B", 9.5)
                pdf.set_xy(MX + 0.10, y + 0.045)
                pdf.cell(CW - 1.2, 0.17, self._pdf_safe(label.upper()), align="L")
                if status:
                    fill, txt, lab = status_style(status)
                    pdf.set_fill_color(*fill)
                    pdf.rect(W - MX - 0.80, y + 0.035, 0.80, 0.19, 'F')
                    pdf.set_text_color(*txt)
                    pdf.set_font("helvetica", "B", 7.5)
                    pdf.set_xy(W - MX - 0.80, y + 0.055)
                    pdf.cell(0.80, 0.15, lab, align="C")
                return y + 0.31

            def cell_lines(text, w, size=7.5):
                pdf.set_font("helvetica", "", size)
                width = pdf.get_string_width(text) if text else 0.0
                return max(1, math.ceil(width * 1.08 / max(0.1, w - 0.12)))

            def table_cell(x, y, w, h, text, size=7.5, bold=False, fill=None, text_c=None, align="L"):
                pdf.set_fill_color(*(fill or (255, 255, 255)))
                pdf.set_draw_color(*LINEC)
                pdf.set_line_width(0.008)
                pdf.rect(x, y, w, h, 'DF')
                pdf.set_text_color(*(text_c or INK))
                pdf.set_font("helvetica", "B" if bold else "", size)
                pdf.set_xy(x + 0.06, y + 0.045)
                pdf.multi_cell(w - 0.12, 0.135, self._pdf_safe(text), align=align)

            # ── Page 1: executive scorecard ──────────────────────────────
            y = new_page()
            y = section_bar(y, "1) Executive Scorecard")
            pdf.set_text_color(*INK)
            pdf.set_font("helvetica", "", 7.5)
            pdf.set_xy(MX, y)
            pdf.cell(CW, 0.14, self._pdf_safe(
                f"This report covers {len(dashboards)} dashboard view(s) exported from Tableau, ordered worst-first. "
                "Each numbered section below contains the cropped dashboard, commentary and a recommended action."))
            y += 0.22

            attention = [d for d in dashboards if d["status"] in ("RED", "AMBER")]
            if attention:
                pdf.set_fill_color(*REDBG)
                pdf.rect(MX, y, CW, 0.24, 'F')
                pdf.set_text_color(*DARKRED)
                pdf.set_font("helvetica", "B", 7.5)
                pdf.set_xy(MX + 0.10, y + 0.05)
                items = "   -   ".join(f"{d['name']} ({d['status']})" for d in attention)
                pdf.cell(CW - 0.2, 0.14, self._pdf_safe(f"NEEDS ATTENTION THIS WEEK:   {items}"))
                y += 0.30

            widths = [0.65, 1.55, 1.85, 3.30, 3.05]
            headers = ["STATUS", "DASHBOARD", "KEY METRIC", "KEY FINDING", "RECOMMENDED ACTION"]
            x = MX
            for wdt, hdr in zip(widths, headers):
                table_cell(x, y, wdt, 0.22, hdr, size=7, bold=True, fill=HDRBG,
                           align="C" if hdr == "STATUS" else "L")
                x += wdt
            y += 0.22

            for i, d in enumerate(dashboards):
                fill = (255, 255, 255) if i % 2 == 0 else ROWALT
                texts = [d["name"], d["key_metric"] or "-", d["headline"], d["recommendation"] or "-"]
                n_lines = max(cell_lines(self._pdf_safe(t), w)
                              for t, w in zip(texts, widths[1:]))
                row_h = n_lines * 0.135 + 0.09
                badge_fill, badge_text, badge_label = status_style(d["status"])
                if d["status"] == "RED":
                    metric_c = RED
                elif d["status"] == "GREEN":
                    metric_c = GREEN
                elif d["status"] == "AMBER":
                    metric_c = (184, 110, 0)
                else:
                    metric_c = INK
                x = MX
                table_cell(x, y, widths[0], row_h, badge_label, bold=True,
                           fill=badge_fill, text_c=badge_text, align="C")
                x += widths[0]
                table_cell(x, y, widths[1], row_h, d["name"], bold=True, fill=fill)
                x += widths[1]
                table_cell(x, y, widths[2], row_h, d["key_metric"] or "-",
                           bold=True, fill=fill, text_c=metric_c)
                x += widths[2]
                table_cell(x, y, widths[3], row_h, d["headline"], fill=fill)
                x += widths[3]
                table_cell(x, y, widths[4], row_h, d["recommendation"] or "-", fill=fill)
                y += row_h
            y += 0.12

            # ── Dashboard sections, flowing down the pages ───────────────
            ZONE_H = 2.15
            for i, d in enumerate(dashboards):
                rec = d["recommendation"]
                act_h = (max(0.26, cell_lines(self._pdf_safe(rec), CW - 0.80) * 0.135 + 0.09)
                         if rec else 0.0)
                needed = 0.31 + 0.30 + ZONE_H + 0.05 + act_h
                if y + needed > BOTTOM:
                    y = new_page()
                y = section_bar(y, f"{i + 2}) {d['name']} - Results", status=d["status"] or None)

                # Yellow HEADLINE strip
                pdf.set_fill_color(*YELLOW)
                pdf.rect(MX, y, CW, 0.24, 'F')
                pdf.set_text_color(*INK)
                pdf.set_font("helvetica", "B", 8.5)
                pdf.set_xy(MX + 0.10, y + 0.05)
                pdf.cell(CW - 0.2, 0.14, self._pdf_safe(f"HEADLINE:  {d['headline']}."))
                y += 0.30

                # Dashboard image (left) with thin grey frame
                try:
                    with Image.open(d["image_path"]) as img_obj:
                        iw, ih = img_obj.size
                    aspect = ih / iw
                    img_w = 4.40
                    img_h = img_w * aspect
                    if img_h > ZONE_H:
                        img_h = ZONE_H
                        img_w = img_h / aspect
                    iy = y + (ZONE_H - img_h) / 2
                    pdf.set_draw_color(*LINEC)
                    pdf.set_line_width(0.01)
                    pdf.rect(MX, iy - 0.02, img_w + 0.04, img_h + 0.04)
                    pdf.image(d["image_path"], x=MX + 0.02, y=iy, w=img_w, h=img_h)
                except Exception as img_err:
                    logging.warning(f"Could not add image to PDF: {img_err}")

                # Comments column (right)
                cmt_x = MX + 4.60
                cmt_w = W - MX - cmt_x
                table_cell(cmt_x, y, cmt_w, 0.20, "COMMENTS", size=7, bold=True, fill=HDRBG)
                old_l, old_r = pdf.l_margin, pdf.r_margin
                pdf.set_left_margin(cmt_x)
                pdf.set_right_margin(W - (cmt_x + cmt_w))
                pdf.set_xy(cmt_x, y + 0.26)
                pdf.set_text_color(*INK)
                findings = d["findings"] or [("Comments", "AI commentary unavailable for this view.")]
                for label, body in findings[:4]:
                    pdf.set_font("helvetica", "B", 7.5)
                    pdf.write(0.125, self._pdf_safe(f"{label}: "))
                    pdf.set_font("helvetica", "", 7.5)
                    pdf.write(0.125, self._pdf_safe(body) + "\n")
                    pdf.ln(0.05)
                meta_bits = []
                ds_text = self._format_datasource_info(d["datasources"])
                filter_text = self._format_filters(d["filters"])
                if ds_text:
                    meta_bits.append(f"Source: {ds_text}")
                if filter_text:
                    meta_bits.append(f"Filters: {filter_text}")
                if meta_bits:
                    pdf.set_text_color(*MUTED)
                    pdf.set_font("helvetica", "I", 6.5)
                    pdf.write(0.11, self._pdf_safe("  |  ".join(meta_bits)))
                pdf.set_left_margin(old_l)
                pdf.set_right_margin(old_r)
                y += ZONE_H + 0.05

                # ACTION row
                if rec:
                    table_cell(MX, y, 0.80, act_h, "ACTION", bold=True,
                               fill=BAR, text_c=(255, 255, 255), align="C")
                    table_cell(MX + 0.80, y, CW - 0.80, act_h, rec, fill=ROWALT)
                    y += act_h
                y += 0.12

            pdf.output(output_path)
            logging.info(f"Successfully created executive WBR-format PDF: {output_path}")
            return output_path

        except Exception as e:
            logging.error(f"Failed to create detailed PDF: {str(e)}", exc_info=True)
            return self.combine_to_pdf(image_paths, output_dir, filename)

    # ── Blend360 theme palette (matched from reference slide) ────────────────
    # Background: very dark navy #060C18
    # Primary teal: #00C4CC  (KPI numbers, borders, rule lines)
    # White: #FFFFFF          (main title, logo)
    # Slate: #94A3B8          (body text, labels)
    # Card bg: #0D1528        (panel fill)
    # Footer bg: #040A12      (darker strip)
    _B360_BG        = RGBColor(0x06, 0x0C, 0x18) if PPTX_AVAILABLE else None
    _B360_TEAL      = RGBColor(0x00, 0xC4, 0xCC) if PPTX_AVAILABLE else None
    _B360_TEAL_DIM  = RGBColor(0x00, 0x6E, 0x76) if PPTX_AVAILABLE else None
    _B360_WHITE     = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
    _B360_LIGHT     = RGBColor(0xB0, 0xBE, 0xCE) if PPTX_AVAILABLE else None
    _B360_CARD      = RGBColor(0x0D, 0x15, 0x28) if PPTX_AVAILABLE else None
    _B360_FOOTER    = RGBColor(0x04, 0x0A, 0x12) if PPTX_AVAILABLE else None

    @staticmethod
    def _b360_bg_xml(slide):
        """Set slide background to Blend360 dark navy."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x06, 0x0C, 0x18)

    @staticmethod
    def _b360_add_rect(slide, x, y, w, h, rgb):
        """Solid filled rectangle with no border line."""
        sh = slide.shapes.add_shape(1, PPTInches(x), PPTInches(y), PPTInches(w), PPTInches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    @staticmethod
    def _b360_add_box(slide, x, y, w, h, fill_rgb, border_rgb, border_pt=1.5):
        """Rectangle with solid fill AND a coloured border line."""
        from pptx.util import Pt as _Pt
        sh = slide.shapes.add_shape(1, PPTInches(x), PPTInches(y), PPTInches(w), PPTInches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_rgb
        sh.line.color.rgb = border_rgb
        sh.line.width = _Pt(border_pt)
        sh.shadow.inherit = False
        return sh

    @staticmethod
    def _b360_text(slide, x, y, w, h, text, size, rgb,
                   bold=False, italic=False, align=PP_ALIGN.LEFT,
                   face="Calibri", wrap=True, margin_in=0.0):
        """Add text box. Returns (txb, tf, first_paragraph)."""
        txb = slide.shapes.add_textbox(PPTInches(x), PPTInches(y), PPTInches(w), PPTInches(h))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        m = Emu(int(margin_in * 914400))
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = m
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = PPTPt(size)
        run.font.color.rgb = rgb
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.name      = face
        return txb, tf, p

    def _b360_header(self, slide, slide_w=13.33, dashboard_name=""):
        """Standard Blend360 header: thin top teal bar, logo, right label, rule."""
        # Very thin teal line at top edge (matches reference)
        self._b360_add_rect(slide, 0, 0, slide_w, 0.045, self._B360_TEAL)
        # BLEND360 logo text
        self._b360_text(slide, 0.38, 0.10, 3.8, 0.42,
                        "BLEND360", 17, self._B360_WHITE,
                        bold=True, face="Arial Black")
        # Right label
        self._b360_text(slide, slide_w - 4.5, 0.15, 4.3, 0.30,
                        "DATA & ANALYTICS PRACTICE", 8, self._B360_TEAL,
                        align=PP_ALIGN.RIGHT, face="Arial")
        # Horizontal teal rule below header
        self._b360_add_rect(slide, 0.38, 0.60, slide_w - 0.76, 0.025, self._B360_TEAL)
        # Dashboard context line (content slides only)
        if dashboard_name:
            self._b360_text(slide, 0.55, 0.68, 10.0, 0.28,
                            f"DATA & ANALYTICS  —  {dashboard_name.upper()}",
                            7.5, self._B360_TEAL, face="Arial")

    def _b360_footer(self, slide, slide_w=13.33, slide_h=7.5, page_num=None):
        """Blend360 footer: dark strip, thin teal rule, confidentiality text, page number."""
        fh = 0.30
        self._b360_add_rect(slide, 0, slide_h - fh, slide_w, fh, self._B360_FOOTER)
        self._b360_add_rect(slide, 0, slide_h - fh, slide_w, 0.022, self._B360_TEAL_DIM)
        self._b360_text(slide, 0.38, slide_h - fh + 0.04, 8, 0.22,
                        "Confidential  |  Blend360  |  Data & Analytics Practice",
                        7.5, self._B360_LIGHT)
        if page_num is not None:
            self._b360_text(slide, slide_w - 1.2, slide_h - fh + 0.04, 1.0, 0.22,
                            str(page_num), 8, self._B360_LIGHT,
                            align=PP_ALIGN.RIGHT)

    # ── Executive Weekly-Business-Review deck format ─────────────────────────
    # Styled after a classic weekly commercial meeting report: white pages,
    # centred brand band, numbered black section bars, yellow HEADLINE strips,
    # plain tables and a narrative Comments column. Sections FLOW down the page
    # like a printed report (multiple sections per page), not one-topic-per-slide.
    _WBR_INK    = RGBColor(0x1A, 0x1A, 0x1A) if PPTX_AVAILABLE else None
    _WBR_BAR    = RGBColor(0x11, 0x11, 0x11) if PPTX_AVAILABLE else None
    _WBR_BAND   = RGBColor(0xEF, 0xEF, 0xEF) if PPTX_AVAILABLE else None
    _WBR_HDRBG  = RGBColor(0xD9, 0xD9, 0xD9) if PPTX_AVAILABLE else None
    _WBR_LINE   = RGBColor(0xAF, 0xAF, 0xAF) if PPTX_AVAILABLE else None
    _WBR_MUTED  = RGBColor(0x59, 0x59, 0x59) if PPTX_AVAILABLE else None
    _WBR_YELLOW = RGBColor(0xFF, 0xF2, 0x00) if PPTX_AVAILABLE else None
    _WBR_ROWALT = RGBColor(0xF5, 0xF5, 0xF5) if PPTX_AVAILABLE else None
    _WBR_WHITE  = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
    _WBR_RED    = RGBColor(0xC0, 0x00, 0x00) if PPTX_AVAILABLE else None
    _WBR_AMBER  = RGBColor(0xE8, 0x9C, 0x00) if PPTX_AVAILABLE else None
    _WBR_GREEN  = RGBColor(0x2E, 0x7D, 0x32) if PPTX_AVAILABLE else None
    _WBR_REDBG  = RGBColor(0xFB, 0xE3, 0xE3) if PPTX_AVAILABLE else None
    _WBR_DARKRED = RGBColor(0x8B, 0x00, 0x00) if PPTX_AVAILABLE else None

    _WBR_TOP    = 0.84   # first content y on a page (below brand band)
    _WBR_BOTTOM = 7.14   # last usable y (above footer)
    _WBR_SEC_H  = 3.13   # uniform height of one dashboard section

    @staticmethod
    def _parse_insight_sections(insights: List[str]):
        """Split raw AI insight strings into (headline, [(LABEL, body), ...])."""
        headline = ""
        sections = []
        if insights:
            headline = insights[0].replace("**", "").strip().rstrip(".")
            for raw in insights[1:]:
                clean = raw.replace("**", "").strip()
                if not clean:
                    continue
                label, sep, body = clean.partition(":")
                if sep and len(label.strip()) <= 25:
                    sections.append((label.strip(), body.strip()))
                else:
                    sections.append(("Insight", clean))
        return headline, sections

    def _wbr_new_page(self, prs, layout, slide_w, slide_h, page_no):
        """Fresh report page: white bg, centred brand band, footer. Returns slide."""
        slide = prs.slides.add_slide(layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._WBR_WHITE

        self._b360_add_rect(slide, 0, 0, slide_w, 0.72, self._WBR_BAND)
        self._b360_add_rect(slide, 0, 0.72, slide_w, 0.014, self._WBR_LINE)
        self._b360_text(slide, 0, 0.08, slide_w, 0.30,
                        "BLEND360", 15, self._WBR_INK,
                        bold=True, align=PP_ALIGN.CENTER, face="Arial Black")
        self._b360_text(slide, 0, 0.40, slide_w, 0.24,
                        "Dashboard Insights — Automated Report to Business Stakeholders",
                        8.5, self._WBR_MUTED, align=PP_ALIGN.CENTER, face="Calibri")

        self._b360_add_rect(slide, 0.30, slide_h - 0.34, slide_w - 0.60, 0.012, self._WBR_LINE)
        self._b360_text(slide, 0.30, slide_h - 0.30, 7.0, 0.22,
                        datetime.now().strftime("Confidential  |  Generated %B %d, %Y"),
                        7.5, self._WBR_MUTED, face="Calibri")
        self._b360_text(slide, slide_w - 1.30, slide_h - 0.30, 1.0, 0.22,
                        f"Page {page_no}", 7.5, self._WBR_MUTED,
                        align=PP_ALIGN.RIGHT, face="Calibri")
        return slide

    def _wbr_section_bar(self, slide, y, slide_w, label):
        """Numbered black section bar, full content width."""
        bar = self._b360_add_rect(slide, 0.30, y, slide_w - 0.60, 0.30, self._WBR_BAR)
        tf = bar.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(int(0.10 * 914400))
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = label.upper()
        run.font.size = PPTPt(10.5)
        run.font.bold = True
        run.font.color.rgb = self._WBR_WHITE
        run.font.name = "Arial"

    def _wbr_cell(self, slide, x, y, w, h, text, size=9.5, bold=False,
                  fill_rgb=None, text_rgb=None, align=PP_ALIGN.LEFT):
        """Bordered table cell with vertically-centred text."""
        box = self._b360_add_box(slide, x, y, w, h,
                                 fill_rgb if fill_rgb is not None else self._WBR_WHITE,
                                 self._WBR_LINE, border_pt=0.75)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(int(0.06 * 914400))
        tf.margin_top = tf.margin_bottom = Emu(int(0.02 * 914400))
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = PPTPt(size)
        run.font.bold = bold
        run.font.color.rgb = text_rgb if text_rgb is not None else self._WBR_INK
        run.font.name = "Calibri"
        return box

    def _wbr_status_style(self, status):
        """(fill, text_colour, label) for a RAG status badge."""
        if status == "RED":
            return self._WBR_RED, self._WBR_WHITE, "RED"
        if status == "AMBER":
            return self._WBR_AMBER, self._WBR_INK, "AMBER"
        if status == "GREEN":
            return self._WBR_GREEN, self._WBR_WHITE, "GREEN"
        return self._WBR_HDRBG, self._WBR_MUTED, "—"

    def _wbr_metric_colour(self, status):
        """Key-metric text colour follows the status semantics."""
        if status == "RED":
            return self._WBR_RED
        if status == "GREEN":
            return self._WBR_GREEN
        if status == "AMBER":
            return RGBColor(0xB8, 0x6E, 0x00)
        return self._WBR_INK

    def _wbr_overview_height(self, n_rows, has_attention):
        """Vertical space the overview section needs."""
        return 0.36 + 0.26 + (0.34 if has_attention else 0.0) + 0.28 + 0.72 * n_rows + 0.12

    def _wbr_overview_section(self, slide, y, slide_w, dashboards):
        """Section 1: executive scorecard — attention strip + RAG status table,
        worst dashboards first. Returns height consumed."""
        self._wbr_section_bar(slide, y, slide_w, "1) Executive Scorecard")
        cy = y + 0.36
        self._b360_text(
            slide, 0.30, cy, slide_w - 0.60, 0.22,
            f"This report covers {len(dashboards)} dashboard view(s) exported from Tableau, ordered worst-first. "
            "Each numbered section below contains the cropped dashboard, commentary and a recommended action.",
            8.5, self._WBR_INK, face="Calibri")
        cy += 0.26

        # "Needs attention" strip — only when something is red/amber
        attention = [d for d in dashboards if d.get("status") in ("RED", "AMBER")]
        if attention:
            self._b360_add_rect(slide, 0.30, cy, slide_w - 0.60, 0.28, self._WBR_REDBG)
            items = "   ·   ".join(f"{d['name']} ({d['status']})" for d in attention)
            self._b360_text(slide, 0.40, cy + 0.03, slide_w - 0.80, 0.22,
                            f"NEEDS ATTENTION THIS WEEK:   {items}",
                            8.5, self._WBR_DARKRED, bold=True, face="Calibri")
            cy += 0.34

        x0 = 0.30
        widths = [0.80, 1.90, 2.20, 3.95, 3.88]
        headers = ["STATUS", "DASHBOARD", "KEY METRIC", "KEY FINDING", "RECOMMENDED ACTION"]
        x = x0
        for wdt, hdr in zip(widths, headers):
            self._wbr_cell(slide, x, cy, wdt, 0.28, hdr, size=8.5, bold=True,
                           fill_rgb=self._WBR_HDRBG,
                           align=PP_ALIGN.CENTER if hdr == "STATUS" else PP_ALIGN.LEFT)
            x += wdt
        cy += 0.28

        for i, d in enumerate(dashboards):
            fill = self._WBR_WHITE if i % 2 == 0 else self._WBR_ROWALT
            badge_fill, badge_text, badge_label = self._wbr_status_style(d.get("status"))
            x = x0
            self._wbr_cell(slide, x, cy, widths[0], 0.72, badge_label,
                           size=9, bold=True, fill_rgb=badge_fill,
                           text_rgb=badge_text, align=PP_ALIGN.CENTER)
            x += widths[0]
            self._wbr_cell(slide, x, cy, widths[1], 0.72, d["name"],
                           size=8.5, bold=True, fill_rgb=fill)
            x += widths[1]
            self._wbr_cell(slide, x, cy, widths[2], 0.72, d.get("key_metric") or "—",
                           size=8.5, bold=True, fill_rgb=fill,
                           text_rgb=self._wbr_metric_colour(d.get("status")))
            x += widths[2]
            self._wbr_cell(slide, x, cy, widths[3], 0.72, d["headline"],
                           size=8.5, fill_rgb=fill)
            x += widths[3]
            self._wbr_cell(slide, x, cy, widths[4], 0.72, d["recommendation"] or "—",
                           size=8.5, fill_rgb=fill)
            cy += 0.72

        return self._wbr_overview_height(len(dashboards), bool(attention))

    def _wbr_dashboard_section(self, slide, y, slide_w, d, sec_no):
        """One numbered dashboard section rendered AT y on the given page:
        section bar, yellow HEADLINE strip, image left + comments right,
        ACTION row. Fixed height _WBR_SEC_H so two sections stack per page."""
        self._wbr_section_bar(slide, y, slide_w, f"{sec_no}) {d['name']} — Results")
        # RAG status chip pinned to the right end of the section bar
        if d.get("status"):
            badge_fill, badge_text, badge_label = self._wbr_status_style(d["status"])
            chip = self._b360_add_rect(slide, slide_w - 0.30 - 0.95, y + 0.045, 0.95, 0.21, badge_fill)
            tf_c = chip.text_frame
            tf_c.word_wrap = False
            tf_c.margin_left = tf_c.margin_right = Emu(0)
            tf_c.margin_top = tf_c.margin_bottom = Emu(0)
            p_c = tf_c.paragraphs[0]
            p_c.alignment = PP_ALIGN.CENTER
            run_c = p_c.add_run()
            run_c.text = badge_label
            run_c.font.size = PPTPt(8)
            run_c.font.bold = True
            run_c.font.color.rgb = badge_text
            run_c.font.name = "Arial"
        cy = y + 0.36

        # Yellow HEADLINE strip
        self._b360_add_rect(slide, 0.30, cy, slide_w - 0.60, 0.28, self._WBR_YELLOW)
        self._b360_text(slide, 0.40, cy + 0.025, slide_w - 0.80, 0.24,
                        f"HEADLINE:  {d['headline']}.",
                        10, self._WBR_INK, bold=True, face="Calibri")
        cy += 0.34

        # Content zone: image (left) + comments (right)
        ZONE_H = 2.00
        IMG_X, IMG_MAX_W = 0.30, 5.40
        try:
            with Image.open(d["image_path"]) as img_obj:
                iw, ih = img_obj.size
            aspect = ih / iw
            img_w = IMG_MAX_W
            img_h = img_w * aspect
            if img_h > ZONE_H:
                img_h = ZONE_H
                img_w = img_h / aspect
            img_x = IMG_X
            img_y = cy + (ZONE_H - img_h) / 2
            self._b360_add_box(slide, img_x - 0.03, img_y - 0.03,
                               img_w + 0.06, img_h + 0.06,
                               self._WBR_WHITE, self._WBR_LINE, border_pt=0.75)
            slide.shapes.add_picture(
                d["image_path"], PPTInches(img_x), PPTInches(img_y),
                width=PPTInches(img_w), height=PPTInches(img_h))
        except Exception as img_err:
            logging.warning(f"Could not add image to slide: {img_err}")

        CMT_X = 5.95
        CMT_W = slide_w - CMT_X - 0.30
        self._wbr_cell(slide, CMT_X, cy, CMT_W, 0.22, "COMMENTS",
                       size=8, bold=True, fill_rgb=self._WBR_HDRBG)

        txb = slide.shapes.add_textbox(
            PPTInches(CMT_X + 0.02), PPTInches(cy + 0.28),
            PPTInches(CMT_W - 0.04), PPTInches(ZONE_H - 0.30))
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        findings = d["findings"] or [("Comments", "AI commentary unavailable for this view.")]
        for label, body in findings[:4]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            lr = p.add_run()
            lr.text = f"{label}: "
            lr.font.size = PPTPt(8)
            lr.font.bold = True
            lr.font.color.rgb = self._WBR_INK
            lr.font.name = "Calibri"
            br = p.add_run()
            br.text = body
            br.font.size = PPTPt(8)
            br.font.color.rgb = self._WBR_INK
            br.font.name = "Calibri"
            p.space_after = PPTPt(4)

        meta_bits = []
        ds_text = self._format_datasource_info(d["datasources"])
        filter_text = self._format_filters(d["filters"])
        if ds_text:
            meta_bits.append(f"Source: {ds_text}")
        if filter_text:
            meta_bits.append(f"Filters: {filter_text}")
        if meta_bits:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            run = p.add_run()
            run.text = "  |  ".join(meta_bits)
            run.font.size = PPTPt(7)
            run.font.italic = True
            run.font.color.rgb = self._WBR_MUTED
            run.font.name = "Calibri"

        cy += ZONE_H + 0.05

        # ACTION row
        if d["recommendation"]:
            self._wbr_cell(slide, 0.30, cy, 0.90, 0.38, "ACTION",
                           size=8.5, bold=True, fill_rgb=self._WBR_BAR,
                           text_rgb=self._WBR_WHITE, align=PP_ALIGN.CENTER)
            self._wbr_cell(slide, 1.20, cy, slide_w - 1.20 - 0.30, 0.38,
                           d["recommendation"], size=8.5, fill_rgb=self._WBR_ROWALT)

        return self._WBR_SEC_H

    def combine_to_pptx_with_details(self, image_paths: List[str], output_dir: str, filename: str, summary_data: List[Dict]) -> str:
        """Build an executive weekly-business-review style PPTX laid out like a
        printed report: numbered sections flow down each page (overview table
        first, then dashboard sections), continuing onto the next page."""
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx is not installed")

        try:
            output_path = os.path.join(output_dir, f"{filename}.pptx")
            prs = Presentation()
            SLIDE_W, SLIDE_H = 13.33, 7.5
            prs.slide_width = PPTInches(SLIDE_W)
            prs.slide_height = PPTInches(SLIDE_H)
            blank_layout = prs.slide_layouts[6]  # fully blank

            # Generate + parse insights for every dashboard up front (worst-first)
            dashboards = self._build_report_dashboards(image_paths, summary_data)

            # ── Flow sections down the pages like a printed report ────────
            page_no = 1
            slide = self._wbr_new_page(prs, blank_layout, SLIDE_W, SLIDE_H, page_no)
            y = self._WBR_TOP

            y += self._wbr_overview_section(slide, y, SLIDE_W, dashboards) + 0.10

            for i, d in enumerate(dashboards):
                if y + self._WBR_SEC_H > self._WBR_BOTTOM:
                    page_no += 1
                    slide = self._wbr_new_page(prs, blank_layout, SLIDE_W, SLIDE_H, page_no)
                    y = self._WBR_TOP
                y += self._wbr_dashboard_section(slide, y, SLIDE_W, d, i + 2) + 0.10

            prs.save(output_path)
            logging.info(f"Successfully created executive WBR-format PPTX: {output_path}")
            return output_path

        except Exception as e:
            logging.error(f"Failed to create PPTX: {str(e)}", exc_info=True)
            raise Exception(f"PPTX creation failed: {str(e)}")

    def combine_to_word(self, image_paths: List[str], output_dir: str, filename: str) -> str:
        """Combine multiple images into a single Word document"""
        try:
            output_path = os.path.join(output_dir, f"{filename}.docx")
            
            # Create new Word document
            doc = Document()
            doc.add_heading('Tableau Dashboard Export', 0)
            
            for i, image_path in enumerate(image_paths):
                if not os.path.exists(image_path):
                    logging.warning(f"Image not found: {image_path}")
                    continue
                
                # Add section heading
                doc.add_heading(f'Dashboard {i + 1}', level=1)
                
                # Add image to document
                # Calculate appropriate width (max 6 inches)
                image = Image.open(image_path)
                aspect_ratio = image.height / image.width
                width = min(6.0, image.width / 100)  # Convert pixels to inches roughly
                height = width * aspect_ratio
                
                doc.add_picture(image_path, width=Inches(width))
                
                # Add page break if not the last image
                if i < len(image_paths) - 1:
                    doc.add_page_break()
            
            # Save document
            doc.save(output_path)
            
            logging.info(f"Successfully created Word document: {output_path}")
            return output_path
            
        except Exception as e:
            logging.error(f"Failed to combine images to Word: {str(e)}")
            raise Exception(f"Word document creation failed: {str(e)}")
    
    def combine_to_word_with_details(self, image_paths: List[str], output_dir: str, filename: str, summary_data: List[Dict]) -> str:
        """Combine multiple images into a single Word document with detailed metadata using 2-column layout"""
        try:
            output_path = os.path.join(output_dir, f"{filename}.docx")
            
            # Create new Word document
            doc = Document()
            
            # Add main title
            title = doc.add_heading('Tableau Dashboard Export Report', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Export Summary removed as per request
            
            # Process dashboards in pairs for same-page layout
            for i in range(0, len(image_paths), 2):
                # Add page heading for dashboard(s)
                # Headers removed as per request
                # if i + 1 < len(image_paths):
                #     page_title = f'Dashboards {i + 1} & {i + 2}'
                # else:
                #     page_title = f'Dashboard {i + 1}'
                
                # doc.add_heading(page_title, level=1)
                
                # First dashboard
                self._add_dashboard_to_word(doc, image_paths[i], summary_data[i], i + 1)
                
                # Second dashboard on same page (if exists)
                if i + 1 < len(image_paths):
                    # Add some spacing between dashboards
                    doc.add_paragraph()
                    self._add_dashboard_to_word(doc, image_paths[i + 1], summary_data[i + 1], i + 2)
                
                # Add page break if not the last pair
                if i + 2 < len(image_paths):
                    doc.add_page_break()
            
            # Save document
            doc.save(output_path)
            
            logging.info(f"Successfully created detailed Word document: {output_path}")
            return output_path
            
        except Exception as e:
            logging.error(f"Failed to combine images to Word with details: {str(e)}")
            raise Exception(f"Detailed Word document creation failed: {str(e)}")
    
    def _add_dashboard_to_word(self, doc, image_path: str, data: Dict, section_num: int):
        """Add a single dashboard to Word document with 2-column layout"""
        try:
            if not os.path.exists(image_path):
                logging.warning(f"Image not found: {image_path}")
                return
            
            # Create a table for 2-column layout
            table = doc.add_table(rows=1, cols=2)
            # Remove table borders
            # table.style = 'Table Grid'
            
            # Set column widths (50% each)
            for cell in table.rows[0].cells:
                cell.width = Inches(3.25)  # Half of 6.5 inch page width
            
            # Left column - Image
            left_cell = table.rows[0].cells[0]
            left_para = left_cell.paragraphs[0]
            
            # Add image to left cell
            image = Image.open(image_path)
            img_width = image.width
            img_height = image.height
            aspect_ratio = img_height / img_width
            
            # Set image width to fit in left column (3 inches max)
            img_width_inches = 3.0
            
            run = left_para.runs[0] if left_para.runs else left_para.add_run()
            run.add_picture(image_path, width=Inches(img_width_inches))
            
            # Right column - Metadata
            right_cell = table.rows[0].cells[1]
            right_para = right_cell.paragraphs[0]
            
            # Metadata and headers removed as per request. Only showing AI insights.
            
            # Add Applied Filters metadata
            applied_filters = data.get("applied_filters", {})
            filter_text = self._format_filters(applied_filters)
            
            if filter_text:
                p_filt = right_cell.add_paragraph()
                run_f_label = p_filt.add_run("Applied Filters: ")
                run_f_label.bold = True
                run_f_label.font.size = Pt(10)
                run_f_val = p_filt.add_run(filter_text)
                run_f_val.font.size = Pt(10)
                p_filt.paragraph_format.space_after = Pt(6)

            # Add Data Source info (As of Date)
            datasources = data.get("datasources", [])
            ds_text = self._format_datasource_info(datasources)
            if ds_text:
                p_ds = right_cell.add_paragraph()
                run_ds_label = p_ds.add_run("Data Source: ")
                run_ds_label.bold = True
                run_ds_label.font.size = Pt(10)
                run_ds_val = p_ds.add_run(ds_text)
                run_ds_val.font.size = Pt(10)
                p_ds.paragraph_format.space_after = Pt(12)

            # Generate and add AI insights if enabled
            if self.gemini_client or self.anthropic_client:
                try:
                    csv_data = data.get('csv_data', '')
                    schema_info = data.get('schema_info', '')
                    logging.info(f"Generating AI insights for dashboard {section_num} (CSV data: {len(csv_data)} chars, Schema info: {len(schema_info)} chars)...")
                    insights = self._generate_ai_insights(image_path, data.get("dashboard", f"Dashboard {section_num}"), csv_data=csv_data, schema_info=schema_info)
                    
                    for idx, insight in enumerate(insights):
                        is_title = (idx == 0)
                        p = right_cell.add_paragraph()
                        
                        if is_title:
                            run = p.add_run(insight)
                            run.bold = True
                            run.font.size = Pt(14)
                            p.paragraph_format.space_after = Pt(12)
                        else:
                            # Add bullet for non-title items
                            p.add_run('• ').bold = True
                            
                            if '**' in insight:
                                parts = insight.split('**')
                                for i, part in enumerate(parts):
                                    run = p.add_run(part)
                                    if i % 2 == 1: # Odd parts were between **
                                        run.bold = True
                            else:
                                p.add_run(insight)
                            
                            # Standard padding
                            p.paragraph_format.space_after = Pt(8)
                    
                    logging.info(f"Successfully added {len(insights)} insights to dashboard {section_num}")
                except Exception as ai_error:
                    logging.warning(f"Failed to generate AI insights for dashboard {section_num}: {ai_error}")
                    # Continue without insights - don't fail the entire document
            
            # Set vertical alignment for cells
            left_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            right_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            
            logging.info(f"Added dashboard {section_num} to Word document in 2-column layout")
            
        except Exception as img_error:
            logging.error(f"Failed to add dashboard {section_num}: {str(img_error)}")
            # Add error message instead
            error_para = doc.add_paragraph()
            error_para.add_run(f'[Error loading Dashboard {section_num}: {os.path.basename(image_path)}]').italic = True
    
    def create_thumbnail(self, image_path: str, max_width: int = 200, max_height: int = 120) -> str:
        """Create a thumbnail of an image"""
        import sys
        try:
            image = Image.open(image_path)

            # Calculate thumbnail size maintaining aspect ratio
            image.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)

            # Generate thumbnail filename
            base_name = os.path.splitext(os.path.basename(image_path))[0]
            thumb_path = os.path.join(os.path.dirname(image_path), f"{base_name}_thumb.png")

            # Ensure the target directory exists
            thumb_dir = os.path.dirname(thumb_path)
            if thumb_dir:
                os.makedirs(thumb_dir, exist_ok=True)

            # On Windows, use the \\?\ long-path prefix when the absolute path
            # would exceed MAX_PATH (260 chars) due to a deeply-nested project dir.
            save_path = thumb_path
            if sys.platform == 'win32':
                abs_path = os.path.abspath(thumb_path)
                if len(abs_path) >= 260:
                    save_path = '\\\\?\\' + abs_path

            image.save(save_path, "PNG")

            logging.info(f"Successfully created thumbnail: {thumb_path}")
            return thumb_path

        except Exception as e:
            logging.error(f"Failed to create thumbnail: {str(e)}")
            raise Exception(f"Thumbnail creation failed: {str(e)}")
    
    def cleanup_temp_files(self):
        """Clean up any temporary files created during processing"""
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except:
                pass
        self.temp_files.clear()
    
    # ── Insights disk cache (pre-warm safety net for live demos) ──────────────
    def _insights_cache_dir(self) -> str:
        d = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".insights_cache")
        try:
            os.makedirs(d, exist_ok=True)
        except Exception:
            pass
        return d

    # Bump this whenever the insight PROMPT changes so cached (old-prompt) insights
    # are automatically invalidated instead of being replayed stale.
    INSIGHTS_PROMPT_VERSION = "v3-2026-06-11"

    def _insights_cache_key(self, image_path: str, dashboard_name: str, csv_data: str, schema_info: str) -> str:
        """Stable key from the inputs that define an insight: prompt version + data + schema + name + image bytes."""
        h = hashlib.md5()
        h.update(self.INSIGHTS_PROMPT_VERSION.encode('utf-8'))
        h.update((dashboard_name or '').encode('utf-8', 'ignore'))
        h.update((csv_data or '').encode('utf-8', 'ignore'))
        h.update((schema_info or '').encode('utf-8', 'ignore'))
        try:
            with open(image_path, 'rb') as f:
                h.update(f.read())
        except Exception:
            pass
        return h.hexdigest()

    @staticmethod
    def _looks_like_error(insights: List[str]) -> bool:
        """Never cache placeholder/error output so a failed run can be retried."""
        if not insights:
            return True
        first = (insights[0] or '').lower()
        return any(t in first for t in ('unavailable', 'failed', 'rate limit', 'no ai provider'))

    def _read_insights_cache(self, key: str):
        if not getattr(config, 'ENABLE_INSIGHTS_CACHE', False):
            return None
        path = os.path.join(self._insights_cache_dir(), f"{key}.json")
        try:
            if os.path.exists(path):
                with open(path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                if isinstance(data, list) and data:
                    return data
        except Exception as e:
            logging.warning(f"Insights cache read failed: {e}")
        return None

    def _write_insights_cache(self, key: str, insights: List[str]) -> None:
        if not getattr(config, 'ENABLE_INSIGHTS_CACHE', False):
            return
        path = os.path.join(self._insights_cache_dir(), f"{key}.json")
        try:
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(insights, f, ensure_ascii=False)
        except Exception as e:
            logging.warning(f"Insights cache write failed: {e}")

    def _generate_ai_insights(self, image_path: str, dashboard_name: str, csv_data: str = '', schema_info: str = '') -> List[str]:
        """
        Cached entry point. On a cache hit the insight is returned with NO API call
        (pre-run a report once before a live demo and the on-stage run is offline).
        On a miss it generates fresh and caches genuine results only.
        """
        key = self._insights_cache_key(image_path, dashboard_name, csv_data, schema_info)
        cached = self._read_insights_cache(key)
        if cached is not None:
            logging.info(f"AI insights served from cache for '{dashboard_name}' (no API call).")
            return cached

        insights = self._generate_ai_insights_uncached(image_path, dashboard_name, csv_data, schema_info)
        if not self._looks_like_error(insights):
            self._write_insights_cache(key, insights)
        return insights

    def _generate_ai_insights_uncached(self, image_path: str, dashboard_name: str, csv_data: str = '', schema_info: str = '') -> List[str]:
        """
        Generate executive AI insights combining the cropped dashboard image with the
        underlying CSV data.  When CSV data is present it is the PRIMARY source of
        truth for all numbers; the image provides visual / structural context only.
        """
        # Use configured provider
        if self.anthropic_client and (config.AI_PROVIDER == "anthropic" or not self.gemini_client):
            return self._generate_anthropic_insights(image_path, dashboard_name, csv_data, schema_info)
        
        if not self.gemini_client:
            return ["AI insights generation failed", "No AI provider configured."]
            
        try:
            pil_image = PILImage.open(image_path)
            has_data = bool(csv_data and csv_data.strip())
            has_schema = bool(schema_info and schema_info.strip())

            if has_data:
                schema_block = f"""
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DATA SOURCE / SCHEMA
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{schema_info.strip()}
""" if has_schema else ''

                prompt = f"""You are the Chief Analytics Officer presenting the single most important read on this
dashboard to a sharp, data-literate executive audience. They can already SEE the dashboard — your job is
to tell them what it MEANS. The raw data below is the PRIMARY truth for every number.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
INPUT 1 — RAW DATA (primary source of truth)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{csv_data.strip()}
{schema_block}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
INPUT 2 — DASHBOARD VISUALIZATION (attached image)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Use the image for chart type, axis labels, colour coding, groupings, visual spikes.
Do NOT read numbers from the image — use Input 1 for all figures.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
WHAT MAKES A REAL INSIGHT (non-negotiable)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Interpret, don't describe. Never just restate a number from the data — expose a relationship,
  comparison, or consequence the viewer would NOT get by glancing at the dashboard.
• Use ONLY numbers present in INPUT 1. Never invent, extrapolate, or round beyond what is shown. Every
  figure must be traceable to a specific cell/row.
• Quantify the comparison every time: percentage difference, a multiple ("3x larger"), share of total,
  or the gap between best and worst.
• Be specific and falsifiable. Name the exact metric, segment, or date. Ban vague phrases like
  "performance is strong" or "focus on growth".
• If the data is a single snapshot with no time dimension, do NOT fabricate a trend or anomaly — report
  concentration, imbalance, or share of total instead, and say so. Credibility beats completeness.
• Use the dashboard's own metric names and segment labels exactly as written.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DELIVER FIVE INSIGHTS — one per dimension
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
1. KPI PERFORMANCE — the headline result: totals/growth/averages, top & bottom performers, exact figures.
2. ANOMALY — the ONE value that deviates most from what the rest of the data would predict (NOT simply
   the largest). Quantify the deviation (e.g. >15% off the norm) and give the most plausible cause.
3. TREND — directional momentum (accelerating / decelerating / reversal), quantified. If no time
   dimension, describe distribution/concentration instead and say so.
4. RISK — the single biggest business risk this data exposes. Tie it to a specific metric and threshold,
   and state the consequence if ignored.
5. RECOMMENDATION — one concrete, owned, time-bound action for the next 30 days, tied to the exact data
   point that makes it urgent.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
OUTPUT FORMAT (return EXACTLY this — no preamble, no closing remarks)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
<TITLE>
A specific CLAIM of 6–10 words that states the finding itself, not a topic.
Good: "Revenue grows 3.5% but AOV and quantity quietly erode".
Bad: "Revenue Performance Overview".
</TITLE>

<BULLET_1>
**KPI Performance:** [Headline metric and change with exact figures and one quantified comparison.] ≤ 45 words.
</BULLET_1>

<BULLET_2>
**Anomaly:** [The value that deviates most from expectation — name it, quantify the deviation, give the cause. Not just "the biggest number".] ≤ 45 words.
</BULLET_2>

<BULLET_3>
**Trend:** [Direction and momentum with a number, plus any inflection or reversal. If no time axis, describe concentration instead.] ≤ 45 words.
</BULLET_3>

<BULLET_4>
**Risk:** [The biggest risk this data exposes, tied to a specific metric/threshold, with the consequence if ignored.] ≤ 45 words.
</BULLET_4>

<BULLET_5>
**Recommendation:** [One concrete 30-day action: what to do, who owns it, the metric to watch — tied to the number that makes it urgent.] ≤ 45 words.
</BULLET_5>

<STATUS>
One word: GREEN, AMBER, or RED. Mechanical rule — RED: the headline metric is declining versus plan or
prior period, OR the Risk requires action this week. AMBER: mixed signals or a watch item. GREEN: on or
above expectation with no material risk. If the headline metric's change is negative, you may NOT answer GREEN.
</STATUS>

<KEY_METRIC>
The single headline figure with its change, ≤ 8 words. Example: "$1,201,846 revenue, +3.5% vs prior 21 days".
</KEY_METRIC>

TONE: Confident, precise, board-room ready. Zero filler. Every bullet carries at least one specific
number or named segment, and tells them something they could NOT get just by looking."""

            else:
                # ── IMAGE-ONLY PROMPT (no CSV available) ────────────────────────
                prompt = """You are the Chief Analytics Officer presenting the single most important read on this
Tableau dashboard to a sharp, data-literate executive audience. They can already SEE the dashboard —
your job is to tell them what it MEANS, not what it shows.

WHAT MAKES A REAL INSIGHT (read before writing):
• Interpret, don't describe. Never restate a number the viewer can already read ("Revenue is $1.2M").
  Instead expose a relationship, comparison, or consequence ("Revenue rose 3.5%, but the entire gain
  came from one channel — the rest are flat").
• Ground every claim in a value visible in the image. Never invent or guess a figure. If a number is
  unreadable, refer to it qualitatively ("the leading channel") — do not fabricate.
• Quantify the comparison every time it's possible: percentage difference, a multiple ("3x larger"),
  share of total, or the gap between best and worst.
• Be specific and falsifiable. Name the exact metric, segment, date, or bar. Ban vague phrases like
  "performance is strong", "do better", or "focus on growth".
• If there is no time dimension, do NOT invent a trend — analyse concentration, imbalance, or the gap
  between leaders and laggards instead.

DELIVER FIVE INSIGHTS — one per dimension:
1. KPI PERFORMANCE — the headline result: the most important metric and what changed, with the figure.
2. ANOMALY — the ONE value that deviates most from what the rest of the data would lead you to expect
   (NOT simply the largest number). Quantify how far it deviates and give the most plausible cause.
3. TREND — direction and momentum over time (accelerating / decelerating / reversing). If no time axis,
   describe how value concentrates across segments and say so.
4. RISK — the single biggest business risk this view exposes. Tie it to a specific metric and state the
   consequence if it is ignored.
5. RECOMMENDATION — one concrete, owned, time-bound action for the next 30 days, tied to the exact
   number that makes it urgent.

OUTPUT FORMAT (return EXACTLY this — no preamble, no closing remarks):
<TITLE>
A specific CLAIM of 6–10 words that states the finding itself, not a topic.
Good: "Revenue grows 3.5% but AOV and quantity quietly erode".
Bad: "Revenue Performance Overview".
</TITLE>

<BULLET_1>
**KPI Performance:** [Headline metric and the change, with an exact figure and one quantified comparison.] ≤ 45 words.
</BULLET_1>

<BULLET_2>
**Anomaly:** [The value that deviates most from expectation — name it, quantify the deviation, give the likely cause. Not just "the biggest bar".] ≤ 45 words.
</BULLET_2>

<BULLET_3>
**Trend:** [Direction and momentum with a number. Name any inflection or reversal. If no time axis, describe concentration instead.] ≤ 45 words.
</BULLET_3>

<BULLET_4>
**Risk:** [The biggest risk this view exposes, tied to a specific metric, with the consequence if ignored.] ≤ 45 words.
</BULLET_4>

<BULLET_5>
**Recommendation:** [One concrete 30-day action: what to do, who owns it, and the metric to watch. Tied to the number that makes it urgent.] ≤ 45 words.
</BULLET_5>

<STATUS>
One word: GREEN, AMBER, or RED. Mechanical rule — RED: the headline metric is declining versus plan or
prior period, OR the Risk requires action this week. AMBER: mixed signals or a watch item. GREEN: on or
above expectation with no material risk. If the headline metric's change is negative, you may NOT answer GREEN.
</STATUS>

<KEY_METRIC>
The single headline figure with its change, ≤ 8 words. Example: "$1.2M revenue, +3.5% vs prior period".
</KEY_METRIC>

TONE: Confident, precise, board-room ready. Zero filler. Every bullet carries at least one specific
number, percentage, or named segment, and tells them something they could NOT get just by looking."""

            # Send image AFTER the prompt so the model reads data first, uses
            # the image only for label/layout context.
            
            # Retry logic and Model Fallback for 429/404 errors
            # We prioritize GEMINI_MODEL from config, but will fallback if needed.
            # IMPORTANT: only list models that currently have quota on this key.
            # gemini-2.0-flash / gemini-2.5-pro return 429 "limit: 0" and only waste
            # retry time, so they are intentionally excluded from the fallback chain.
            fallback_models = [GEMINI_MODEL, 'gemini-2.5-flash', 'gemini-2.5-flash-lite']
            # Deduplicate while preserving order
            models_to_try = []
            for m in fallback_models:
                if m not in models_to_try:
                    models_to_try.append(m)
            
            logging.info(f"AI EXECUTION: Attempting insights using model sequence: {models_to_try}")
            print(f"\n[AI DEBUG] Models to try: {models_to_try}")
                    
            max_retries_per_model = 2  # Reduced from 5 — fast fail and move to next model
            retry_delay = 5            # Reduced from 10
            response = None
            last_error = Exception("All Gemini models failed")

            for model_name in models_to_try:
                for attempt in range(max_retries_per_model):
                    try:
                        logging.info(f"Attempting AI insights with model: {model_name} (Attempt {attempt+1}/{max_retries_per_model})")
                        response = self.gemini_client.models.generate_content(
                            model=model_name,
                            contents=[prompt, pil_image]
                        )
                        break # Success with this model!
                    except Exception as e:
                        last_error = e
                        error_msg = str(e).lower()

                        # Handle 404 (Model not found) - try next model immediately
                        if "404" in error_msg or "not_found" in error_msg:
                            logging.warning(f"Model {model_name} not found (404). Trying fallback...")
                            break # Break retry loop to try next model

                        # Handle 429 (rate limit) and 503 (transient high demand) —
                        # both are temporary, so back off and retry the same model
                        # once before moving to the fallback.
                        if ("429" in error_msg or "resource_exhausted" in error_msg
                                or "503" in error_msg or "unavailable" in error_msg or "overloaded" in error_msg):
                            kind = "Rate limit" if ("429" in error_msg or "resource_exhausted" in error_msg) else "Model busy (503)"
                            if attempt < max_retries_per_model - 1:
                                wait_time = retry_delay * (attempt + 1)  # linear backoff
                                msg = f"{kind} for {model_name}. Waiting {wait_time}s... (Attempt {attempt+1}/{max_retries_per_model})"
                                logging.warning(msg)
                                print(f"\n[AI INSIGHTS] {msg}")
                                time.sleep(wait_time)
                                continue
                            else:
                                logging.warning(f"{kind} persisted for {model_name}. Trying fallback...")
                                break # Try next model

                        # Other errors - try next model
                        logging.error(f"Error with {model_name}: {str(e)}. Trying fallback...")
                        break
                
                if response:
                    break # Success!
                    
            if not response:
                error_summary = f"All Gemini models {models_to_try} failed. Last error: {str(last_error)}"
                logging.error(error_summary)
                raise Exception(error_summary)

            raw = response.text.strip()

            # ── Parse structured XML-like tags ────────────────────────────────
            import re
            def extract_tag(tag, text):
                m = re.search(rf'<{tag}>(.*?)</{tag}>', text, re.DOTALL)
                return m.group(1).strip() if m else None

            title   = extract_tag('TITLE',    raw)
            bullet1 = extract_tag('BULLET_1', raw)
            bullet2 = extract_tag('BULLET_2', raw)
            bullet3 = extract_tag('BULLET_3', raw)
            bullet4 = extract_tag('BULLET_4', raw)
            bullet5 = extract_tag('BULLET_5', raw)
            status     = extract_tag('STATUS',     raw)
            key_metric = extract_tag('KEY_METRIC', raw)

            if title and bullet1:
                insights = [title]
                for b in (bullet1, bullet2, bullet3, bullet4, bullet5):
                    if b:
                        insights.append(b)
                if status:
                    s = status.strip().upper()
                    for tok in ("RED", "AMBER", "GREEN"):
                        if tok in s:
                            insights.append(f"**Status:** {tok}")
                            break
                if key_metric:
                    insights.append(f"**Key Metric:** {key_metric.strip()}")
                return insights[:8]

            # ── Fallback: plain-line parsing (handles non-tagged responses) ───
            lines = [l.strip() for l in raw.split('\n') if l.strip()]
            insights = []
            if lines:
                insights.append(lines[0].lstrip('•-*#1234. ').strip())
                for line in lines[1:]:
                    clean = line.lstrip('•-*#1234. ').strip()
                    if clean:
                        insights.append(clean)

            return insights[:4] if insights else ["Insight Summary unavailable"]

        except Exception as e:
            logging.error(f"Failed to generate AI insights: {str(e)}")
            return ["AI insights unavailable", "Check server logs for details."]

    def _generate_anthropic_insights(self, image_path: str, dashboard_name: str, csv_data: str = '', schema_info: str = '') -> List[str]:
        """Generate insights using Anthropic Claude"""
        try:
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            prompt = f"Analyze this dashboard: {dashboard_name}. Data: {csv_data}. Schema: {schema_info}. Provide 5 bullet points (KPI, Anomaly, Trend, Risk, Recommendation)."
            
            message = self.anthropic_client.messages.create(
                model=config.ANTHROPIC_MODEL,
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": image_data}},
                        {"type": "text", "text": prompt}
                    ]
                }]
            )
            return message.content[0].text.split('\n')
        except Exception as e:
            error_msg = str(e)
            logging.error(f"Anthropic insights failed: {error_msg}")
            # Identify specific errors
            if "Authentication" in error_msg or "401" in error_msg:
                return ["Anthropic Authentication Failed", "Please check your API key in config.py."]
            if "Resource" in error_msg or "429" in error_msg:
                return ["Claude Rate Limit Hit", "Your current tier allows 5 requests per minute. Please wait."]
            return ["Anthropic insights generation failed", f"Error: {error_msg[:100]}..."]

    def vision_classifier(self, image_path: str) -> str:
        """AI Vision: Identify the dashboard type and business context"""
        if not self.gemini_client:
            return "Classification Unavailable (AI Disabled)"
        
        try:
            pil_image = PILImage.open(image_path)
            prompt = "Analyze this Tableau dashboard image. Identify the primary industry (e.g., Sales, Finance, Supply Chain) and the intended audience (e.g., Executive, Operational). Return a single concise string like 'Executive Sales Dashboard'."
            
            response = self.gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[prompt, pil_image]
            )
            return response.text.strip()
        except Exception as e:
            logging.error(f"Vision Classifier failed: {e}")
            return "Unknown Dashboard Type"

    def vision_dashboard_analyze(self, image_path: str) -> List[str]:
        """AI Vision: Discover all worksheet names and visual components visible in the image"""
        if not self.gemini_client:
            return []
            
        try:
            pil_image = PILImage.open(image_path)
            prompt = """Analyze the attached Tableau dashboard image.
Extract the names of all individual worksheets, charts, or data tables visible in this dashboard.
Return ONLY a JSON list of strings. No other text.

Example Output:
["Sheet 1", "Product Sales", "Trend Line"]
"""
            response = self.gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[prompt, pil_image]
            )
            
            text = response.text.strip()
            # Try to extract JSON list
            if "[" in text and "]" in text:
                import json
                start = text.find("[")
                end = text.rfind("]") + 1
                try:
                    names = json.loads(text[start:end])
                    return list(set(names))
                except:
                    pass
            
            # Fallback to line-by-line if JSON fails
            lines = [l.strip().lstrip('•-1234. ') for l in text.split('\n') if l.strip()]
            return list(set(lines)) # Deduplicate
        except Exception as e:
            logging.error(f"Vision Dashboard Analyze failed: {e}")
            return []

    def vision_image_segment(self, image_path: str) -> Dict[str, Any]:
        """AI Vision: Segment the image into logical layout areas (Classic CV + AI)"""
        try:
            # Load with OpenCV for segmentation
            img = cv2.imread(image_path)
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            
            # Simple thresholding to find boundaries
            _, thresh = cv2.threshold(gray, 240, 255, cv2.THRESH_BINARY_INV)
            contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            segments = []
            for cnt in contours:
                x, y, w, h = cv2.boundingRect(cnt)
                if w > 50 and h > 50: # Filter noise
                    segments.append({'x': x, 'y': y, 'w': w, 'h': h})
            
            # Use AI to describe the overall layout
            layout_desc = "Standard Grid"
            if self.gemini_client:
                pil_image = PILImage.open(image_path)
                prompt = "Describe the layout of this dashboard in 10 words (e.g., 'Two-column layout with header and three charts below')."
                resp = self.gemini_client.models.generate_content(model=GEMINI_MODEL, contents=[prompt, pil_image])
                layout_desc = resp.text.strip()
                
            return {
                'segment_count': len(segments),
                'layout_description': layout_desc,
                'raw_segments': segments[:10] # Top 10 for info
            }
        except Exception as e:
            logging.error(f"Vision Image Segment failed: {e}")
            return {'error': str(e)}
