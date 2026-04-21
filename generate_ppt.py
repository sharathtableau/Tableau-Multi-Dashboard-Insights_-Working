from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0D, 0x1B, 0x3E)   # deep blue – title backgrounds
TEAL     = RGBColor(0x00, 0x7A, 0x8A)   # accent / dividers
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)   # soft off-white for body slides
DARK_TXT = RGBColor(0x1A, 0x1A, 0x2E)   # near-black for readability
ACCENT   = RGBColor(0x00, 0xC2, 0xCB)   # bright teal for highlights
GREEN    = RGBColor(0x2E, 0x7D, 0x32)


def _bg(slide, color: RGBColor):
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def _add_bullet_box(slide, items, left, top, width, height,
                    font_size=16, title=None, title_size=18,
                    text_color=DARK_TXT, title_color=TEAL):
    """Add a textbox with optional bolded title + bulleted list of (text, level) tuples."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"

    for text, level in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = "Calibri"
    return txb


def _divider(slide, top_inches, color=TEAL, width=9.0, left=0.5):
    """Draw a thin horizontal line as a visual divider."""
    line = slide.shapes.add_connector(
        1,  # straight
        Inches(left), Inches(top_inches),
        Inches(left + width), Inches(top_inches)
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 – Title / Hero"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _bg(slide, NAVY)

    # Product name
    _add_textbox(slide, "Snapshot Insights",
                 0.4, 1.5, 9.2, 1.2,
                 font_size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Tagline
    _add_textbox(slide,
                 "Tableau Dashboard Cropper & AI Insights Generator",
                 0.4, 2.8, 9.2, 0.7,
                 font_size=22, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Sub-tagline
    _add_textbox(slide,
                 "From Dashboard to Decision-Ready Report — in Minutes, Not Hours.",
                 0.4, 3.55, 9.2, 0.6,
                 font_size=16, bold=False,
                 color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    # Author tag
    _add_textbox(slide, "Built by Sharath Kumar Kammari  |  Client Demo Overview",
                 0.4, 6.4, 9.2, 0.4,
                 font_size=11, bold=False,
                 color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


def slide_purpose(prs):
    """Slide 2 – Purpose"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT_BG)

    _add_textbox(slide, "Purpose", 0.5, 0.25, 9.0, 0.65,
                 font_size=30, bold=True, color=NAVY)
    _divider(slide, 1.05)

    _add_textbox(slide,
                 "Organizations using Tableau spend hours manually exporting dashboards, "
                 "screenshotting charts, digging up underlying data, and formatting reports "
                 "for stakeholders — a process that is repetitive, error-prone, and slow.",
                 0.5, 1.20, 9.0, 1.1,
                 font_size=14, bold=False, color=DARK_TXT)

    _add_textbox(slide,
                 "Snapshot Insights was built to eliminate this bottleneck.",
                 0.5, 2.40, 9.0, 0.45,
                 font_size=15, bold=True, color=TEAL)

    bullets = [
        ("Securely connect to any Tableau Online or Tableau Server environment", 0),
        ("Visually select and crop only the metrics that matter", 0),
        ("Automatically extract the exact underlying backend data for that crop", 0),
        ("Generate an AI-powered, insight-rich report — ready for immediate distribution", 0),
    ]
    _add_bullet_box(slide, bullets, 0.5, 2.95, 9.0, 2.6,
                    font_size=14, text_color=DARK_TXT)

    _add_textbox(slide,
                 "One tool · Zero manual steps · Entire reporting pipeline automated.",
                 0.5, 6.1, 9.0, 0.45,
                 font_size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def slide_outcome(prs):
    """Slide 3 – Outcome"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    _add_textbox(slide, "Outcome", 0.5, 0.25, 9.0, 0.65,
                 font_size=30, bold=True, color=NAVY)
    _divider(slide, 1.05)

    _add_textbox(slide,
                 "By the end of a session with Snapshot Insights, a user will have:",
                 0.5, 1.15, 9.0, 0.45,
                 font_size=15, bold=False, color=DARK_TXT)

    outcomes = [
        ("1.  Exported  high-fidelity PDFs of up to 6 Tableau dashboards directly from the server", 0),
        ("2.  Cropped  the exact chart segments they need via a drag-and-drop browser interface", 0),
        ("3.  Received  the precise raw backend crosstab data powering those visuals", 0),
        ("4.  Obtained  a professionally formatted Word / PDF report containing:", 0),
        ("The cropped dashboard image + metadata (Project, Workbook, Timestamp)", 1),
        ("AI-generated business insights synthesised from the visual AND raw data numbers", 1),
    ]
    _add_bullet_box(slide, outcomes, 0.5, 1.70, 9.0, 3.6,
                    font_size=14, text_color=DARK_TXT)

    _add_textbox(slide,
                 "Result: A document an executive can act on immediately — "
                 "no manual formatting, data-gathering, or narrative writing required.",
                 0.5, 5.55, 9.0, 0.75,
                 font_size=14, bold=True, color=NAVY)


def slide_business_impact(prs):
    """Slide 4 – Business Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT_BG)

    _add_textbox(slide, "Business Impact", 0.5, 0.25, 9.0, 0.65,
                 font_size=30, bold=True, color=NAVY)
    _divider(slide, 1.05)

    # --- Before / After comparison table (manual textboxes) ----
    col_headers = ["Impact Area", "Before", "After"]
    col_x       = [0.45, 3.20, 6.45]
    col_w       = [2.70, 3.20, 3.20]

    rows = [
        ("Report Preparation",   "2–4 hours per cycle",                        "5–10 minutes"),
        ("Data Accuracy",        "Manual copy/paste errors",                   "Tied to live backend data"),
        ("Insights Quality",     "Analyst writes narrative by hand",           "AI synthesises visual + data"),
        ("Consistency",          "Format varies by person",                    "Standardised every time"),
        ("Scalability",          "Limited by human bandwidth",                 "6 dashboards simultaneously"),
        ("Tool Dependency",      "Tableau Desktop + Office skills needed",     "Browser only, zero install"),
    ]

    # Header row
    for i, hdr in enumerate(col_headers):
        hdr_color = WHITE if i == 0 else WHITE
        bg_color  = NAVY if i == 0 else TEAL
        box = slide.shapes.add_shape(
            1,  # rectangle
            Inches(col_x[i]), Inches(1.20),
            Inches(col_w[i]), Inches(0.38)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = WHITE
        tf2 = box.text_frame
        tf2.paragraphs[0].text = hdr
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        run2 = tf2.paragraphs[0].runs[0]
        run2.font.bold = True
        run2.font.size = Pt(12)
        run2.font.color.rgb = hdr_color
        run2.font.name = "Calibri"

    # Data rows
    for r, (area, before, after) in enumerate(rows):
        row_top = 1.58 + r * 0.56
        row_bg  = RGBColor(0xE8, 0xF0, 0xFE) if r % 2 == 0 else WHITE
        for c, cell_val in enumerate([area, before, after]):
            box = slide.shapes.add_shape(
                1,
                Inches(col_x[c]), Inches(row_top),
                Inches(col_w[c]), Inches(0.52)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = row_bg
            box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            tf3 = box.text_frame
            tf3.word_wrap = True
            tf3.margin_left  = Inches(0.05)
            tf3.margin_right = Inches(0.05)
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            run3 = p3.add_run()
            run3.text = cell_val
            run3.font.size = Pt(11)
            run3.font.name = "Calibri"
            run3.font.bold = (c == 0)
            run3.font.color.rgb = GREEN if c == 2 else DARK_TXT


def slide_who(prs):
    """Slide 5 – Who Is This For"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    _add_textbox(slide, "Who Is This Most Useful For?", 0.5, 0.25, 9.0, 0.65,
                 font_size=28, bold=True, color=NAVY)
    _divider(slide, 1.05)

    # Left column – Personas
    personas = [
        ("Business Analysts", "Eliminates Export → Screenshot → Paste → Format cycles"),
        ("Data / BI Managers", "Centralise report generation without needing Tableau Desktop"),
        ("Project Managers",   "Extract KPIs from shared dashboards without a BI team"),
        ("Executives / Leaders", "Receive clean, AI-narrated reports — proactively"),
    ]
    _add_textbox(slide, "Primary Users", 0.5, 1.15, 4.4, 0.38,
                 font_size=15, bold=True, color=TEAL)
    bullet_persona = []
    for persona, why in personas:
        bullet_persona.append((f"► {persona}", 0))
        bullet_persona.append((f"   {why}", 1))
    _add_bullet_box(slide, bullet_persona, 0.5, 1.55, 4.4, 4.0,
                    font_size=12, text_color=DARK_TXT)

    # Right column – Industries
    industries = [
        ("Retail & E-Commerce — weekly sales & regional performance reporting", 0),
        ("SaaS / Technology — product metrics, adoption, revenue tracking", 0),
        ("Finance & Consulting — client-ready reports from internal BI platforms", 0),
        ("Healthcare / Operations — standardising periodic performance reviews", 0),
    ]
    _add_textbox(slide, "Best-Fit Industries", 5.1, 1.15, 4.4, 0.38,
                 font_size=15, bold=True, color=TEAL)
    _add_bullet_box(slide, industries, 5.1, 1.55, 4.4, 3.0,
                    font_size=12, text_color=DARK_TXT)

    # Team profile band at bottom
    _add_textbox(slide,
                 "Any team using Tableau Online / Server that prepares regular stakeholder reports "
                 "and wants to reduce manual effort while adding AI-powered narrative.",
                 0.5, 6.05, 9.0, 0.65,
                 font_size=13, bold=False, color=WHITE)
    box = slide.shapes.add_shape(
        1, Inches(0.45), Inches(5.95), Inches(9.1), Inches(0.80)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.zorder = 0   # send behind the textbox


def slide_workflow(prs):
    """Slide 6 – Core Workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT_BG)

    _add_textbox(slide, "How It Works — End-to-End Workflow", 0.5, 0.25, 9.0, 0.65,
                 font_size=28, bold=True, color=NAVY)
    _divider(slide, 1.05)

    steps = [
        ("1", "Secure Authentication",     "Users log in via Tableau Server / Cloud credentials"),
        ("2", "Dashboard Selection",       "Browse Project → Workbook → Dashboard using the Tableau REST API"),
        ("3", "PDF → PNG Conversion",      "High-fidelity PDF exported and rendered as an interactive image"),
        ("4", "Interactive Crop",          "Drag-and-drop bounding box to isolate the exact chart segment"),
        ("5", "Headless Data Extraction",  "Selenium downloads the exact Crosstab Excel behind the visual"),
        ("6", "AI Visual-Data Bridge",     "Gemini Vision maps the crop to the correct data columns"),
        ("7", "Report Generation",         "Word / PDF assembled with image, metadata & AI insight block"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        top = 1.20 + i * 0.68
        # Numbered circle
        circ = slide.shapes.add_shape(
            9,  # oval
            Inches(0.45), Inches(top),
            Inches(0.42), Inches(0.42)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        tf_c = circ.text_frame
        tf_c.paragraphs[0].text = num
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_c = tf_c.paragraphs[0].runs[0]
        run_c.font.bold = True
        run_c.font.size = Pt(12)
        run_c.font.color.rgb = WHITE
        run_c.font.name = "Calibri"
        # Step text
        _add_textbox(slide, f"{title}  —  {desc}",
                     1.0, top - 0.04, 8.6, 0.48,
                     font_size=13, bold=False, color=DARK_TXT)


def slide_tech_stack(prs):
    """Slide 7 – Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    _add_textbox(slide, "Technology Stack", 0.5, 0.25, 9.0, 0.65,
                 font_size=30, bold=True, color=NAVY)
    _divider(slide, 1.05)

    categories = [
        ("Web Framework & Hosting",   "Flask · Gunicorn · Render.com"),
        ("Integration & APIs",        "Tableau REST API · Selenium (headless Chrome) · webdriver-manager"),
        ("AI & Data Processing",      "Google Gemini (Vision + Text) · Pandas"),
        ("Image & PDF Handling",      "Pillow (PIL) · pdf2image (Poppler)"),
        ("Document Generation",       "python-docx (Word) · fpdf2 (PDF) · python-pptx"),
        ("Frontend",                  "HTML · CSS · JavaScript · Bootstrap"),
    ]

    for i, (category, tools) in enumerate(categories):
        top = 1.20 + i * 0.82
        # Category pill
        box = slide.shapes.add_shape(
            1,
            Inches(0.45), Inches(top),
            Inches(2.85), Inches(0.40)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        tf_b = box.text_frame
        tf_b.paragraphs[0].text = category
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_b = tf_b.paragraphs[0].runs[0]
        run_b.font.bold = True
        run_b.font.size = Pt(11)
        run_b.font.color.rgb = WHITE
        run_b.font.name = "Calibri"
        # Tools text
        _add_textbox(slide, tools,
                     3.45, top - 0.02, 6.1, 0.44,
                     font_size=13, bold=False, color=DARK_TXT)


def slide_closing(prs):
    """Slide 8 – Closing / CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)

    _add_textbox(slide, "Thank You", 0.4, 1.4, 9.2, 0.9,
                 font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "Snapshot Insights — Turning Tableau Dashboards into\n"
                 "Decision-Ready Reports Powered by AI",
                 0.4, 2.45, 9.2, 0.9,
                 font_size=18, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    _divider(slide, 3.55, color=ACCENT, width=5.0, left=2.5)

    _add_textbox(slide,
                 "Live Demo:  https://tableaudashboardcropper.onrender.com/login",
                 0.4, 3.75, 9.2, 0.5,
                 font_size=14, bold=False,
                 color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "Built by Sharath Kumar Kammari  |  Open for collaborations & custom deployments",
                 0.4, 6.3, 9.2, 0.4,
                 font_size=11, bold=False,
                 color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def create_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_purpose(prs)
    slide_outcome(prs)
    slide_business_impact(prs)
    slide_who(prs)
    slide_workflow(prs)
    slide_tech_stack(prs)
    slide_closing(prs)

    out = "Snapshot_Insights_Demo.pptx"
    prs.save(out)
    print(f"✅  Presentation saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    create_presentation()
